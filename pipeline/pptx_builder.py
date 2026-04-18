import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}

_DEFAULT_THEME = {
    "bg":      "1A1A2E",
    "surface": "16213E",
    "accent":  "4D9DFF",
    "text":    "FFFFFF",
    "subtext": "AABBCC",
    "body":    "DDEEFF",
}


def _parse_hex(s: str) -> RGBColor:
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _blend(hex1: str, hex2: str, ratio: float) -> RGBColor:
    """Blend hex1 toward hex2 by ratio (0=hex1, 1=hex2)."""
    h1, h2 = hex1.lstrip("#"), hex2.lstrip("#")
    r = int(int(h1[0:2], 16) * (1 - ratio) + int(h2[0:2], 16) * ratio)
    g = int(int(h1[2:4], 16) * (1 - ratio) + int(h2[2:4], 16) * ratio)
    b = int(int(h1[4:6], 16) * (1 - ratio) + int(h2[4:6], 16) * ratio)
    return RGBColor(r, g, b)


def _make_theme(raw: dict) -> dict:
    merged = {**_DEFAULT_THEME}
    for k, v in raw.items():
        if v and k in merged:
            merged[k] = v.lstrip("#")
    return merged


def _rgb(theme: dict, key: str) -> RGBColor:
    return _parse_hex(theme[key])


def build_pptx(slides_data: dict, output_dir: str) -> str:
    theme = _make_theme(slides_data.get("theme", {}))
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    dispatch = {
        "title":      _title_slide,
        "agenda":     _agenda_slide,
        "section":    _section_slide,
        "stat":       _stat_slide,
        "two_column": _two_column_slide,
        "quote":      _quote_slide,
        "chart":      _chart_slide,
        "summary":    _summary_slide,
    }

    for slide in slides_data.get("slides", []):
        fn = dispatch.get(slide.get("type", "content"), _content_slide)
        fn(prs, slide, theme)

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r"[^\w가-힣]", "_", slides_data.get("title", "presentation"))
    out_path = str(Path(output_dir) / f"{safe_name}_{ts}.pptx")
    prs.save(out_path)
    return out_path


# ── helpers ──────────────────────────────────────────────────────────────────

def _blank(prs: Presentation, theme: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme, "bg")
    return slide


def _textbox(slide, left, top, width, height, text, size,
             bold=False, color=None, align=PP_ALIGN.LEFT, theme=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color else _rgb(theme, "text")
    return tf


def _accent_bar(slide, top: float, theme: dict, left=0.8, width=1.4):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(3),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme, "accent")
    bar.line.fill.background()


def _rect(slide, x, y, w, h, fill_color: RGBColor, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ── slide renderers ───────────────────────────────────────────────────────────

def _title_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)
    # Accent band at bottom
    _rect(slide, 0, 6.8, 13.33, 0.7, _rgb(theme, "accent"))
    _textbox(slide, 1, 2.2, 11.33, 1.8, data.get("title", ""), 44,
             bold=True, align=PP_ALIGN.CENTER, theme=theme)
    subtitle = data.get("subtitle", "")
    if subtitle:
        _textbox(slide, 1, 4.1, 11.33, 1.0, subtitle, 24,
                 color=_rgb(theme, "subtext"), align=PP_ALIGN.CENTER, theme=theme)


def _agenda_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)
    _textbox(slide, 0.8, 0.4, 11.5, 0.9, data.get("title", "목차"), 32,
             bold=True, theme=theme)
    _accent_bar(slide, 1.25, theme)

    items = data.get("items", [])
    box = slide.shapes.add_textbox(Inches(2.0), Inches(1.6), Inches(9.5), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(16)
        run_num = p.add_run()
        run_num.text = f"{i + 1:02d}  "
        run_num.font.size = Pt(24)
        run_num.font.bold = True
        run_num.font.color.rgb = _rgb(theme, "accent")
        run_txt = p.add_run()
        run_txt.text = item
        run_txt.font.size = Pt(24)
        run_txt.font.color.rgb = _rgb(theme, "body")


def _section_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)

    # Ghost diamond — blend accent 20% over bg for subtle watermark
    ghost_color = _blend(theme["bg"], theme["accent"], 0.18)
    ghost_box = slide.shapes.add_textbox(Inches(6.5), Inches(-0.5), Inches(7), Inches(8))
    tf_g = ghost_box.text_frame
    p_g = tf_g.paragraphs[0]
    p_g.alignment = PP_ALIGN.RIGHT
    run_g = p_g.add_run()
    run_g.text = "◆"
    run_g.font.size = Pt(340)
    run_g.font.bold = True
    run_g.font.color.rgb = ghost_color

    # Left vertical accent bar
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.35), Inches(1.6), Pt(5), Inches(4.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme, "accent")
    bar.line.fill.background()

    _textbox(slide, 0.9, 2.3, 8.5, 1.8, data.get("title", ""), 48,
             bold=True, theme=theme)
    subtitle = data.get("subtitle", "")
    if subtitle:
        _textbox(slide, 0.9, 4.3, 8.5, 1.0, subtitle, 22,
                 color=_rgb(theme, "subtext"), theme=theme)


def _content_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)
    _textbox(slide, 0.8, 0.4, 11.5, 0.9, data.get("title", ""), 32,
             bold=True, theme=theme)
    _accent_bar(slide, 1.25, theme)

    bullets = data.get("bullets", [])
    if not bullets:
        return
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.5), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:3]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(22)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size = Pt(22)
        run.font.color.rgb = _rgb(theme, "body")


def _stat_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)
    _textbox(slide, 0.8, 0.4, 11.5, 0.9, data.get("title", ""), 32,
             bold=True, theme=theme)
    _accent_bar(slide, 1.25, theme)

    stats = data.get("stats", [])[:3]
    if not stats:
        return

    count = len(stats)
    total_w = 11.0
    card_w = (total_w - (count - 1) * 0.4) / count
    start_x = 1.15

    for i, stat in enumerate(stats):
        x = start_x + i * (card_w + 0.4)
        # Card background
        _rect(slide, x, 1.7, card_w, 4.8,
              _rgb(theme, "surface"), _rgb(theme, "accent"))

        # Top accent stripe on card
        _rect(slide, x, 1.7, card_w, 0.12, _rgb(theme, "accent"))

        # Value
        val_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(2.4), Inches(card_w - 0.2), Inches(2.0))
        tf_v = val_box.text_frame
        p_v = tf_v.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        run_v = p_v.add_run()
        run_v.text = stat.get("value", "")
        run_v.font.size = Pt(54)
        run_v.font.bold = True
        run_v.font.color.rgb = _rgb(theme, "accent")

        # Label
        lbl_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(4.7), Inches(card_w - 0.2), Inches(1.4))
        tf_l = lbl_box.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.alignment = PP_ALIGN.CENTER
        run_l = p_l.add_run()
        run_l.text = stat.get("label", "")
        run_l.font.size = Pt(17)
        run_l.font.color.rgb = _rgb(theme, "subtext")


def _two_column_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)
    _textbox(slide, 0.8, 0.4, 11.5, 0.9, data.get("title", ""), 32,
             bold=True, theme=theme)
    _accent_bar(slide, 1.25, theme)

    # Left panel
    _rect(slide, 0.4, 1.45, 6.0, 5.75, _rgb(theme, "surface"))
    # Right panel
    _rect(slide, 6.9, 1.45, 6.0, 5.75, _rgb(theme, "surface"))
    # Center divider
    _rect(slide, 6.55, 1.6, 0.05, 5.5, _rgb(theme, "accent"))

    for panel_data, px in [(data.get("left", {}), 0.6), (data.get("right", {}), 7.1)]:
        heading = panel_data.get("heading", "")
        bullets = panel_data.get("bullets", [])
        if heading:
            _textbox(slide, px, 1.6, 5.6, 0.75, heading, 20,
                     bold=True, color=_rgb(theme, "accent"), theme=theme)
        if bullets:
            box = slide.shapes.add_textbox(
                Inches(px), Inches(2.45), Inches(5.6), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets[:4]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(14)
                run = p.add_run()
                run.text = f"  {b}"
                run.font.size = Pt(18)
                run.font.color.rgb = _rgb(theme, "body")


def _quote_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)

    # Accent band at top
    _rect(slide, 0, 0, 13.33, 0.12, _rgb(theme, "accent"))
    # Accent band at bottom
    _rect(slide, 0, 7.38, 13.33, 0.12, _rgb(theme, "accent"))

    # Large opening quote mark
    _textbox(slide, 0.5, 0.4, 2.0, 2.5, "\u275d", 100,
             bold=True, color=_rgb(theme, "accent"), theme=theme)

    # Quote text
    quote = data.get("quote", "")
    _textbox(slide, 1.5, 1.6, 10.5, 3.8, quote, 28,
             color=_rgb(theme, "text"), align=PP_ALIGN.LEFT, theme=theme)

    # Attribution
    attribution = data.get("attribution", "")
    if attribution:
        _textbox(slide, 1.5, 5.6, 10.5, 0.9, f"— {attribution}", 18,
                 color=_rgb(theme, "subtext"), theme=theme)


def _summary_slide(prs, data: dict, theme: dict):
    slide = _blank(prs, theme)

    # Header band (accent-colored)
    _rect(slide, 0, 0, 13.33, 1.5, _rgb(theme, "accent"))
    _textbox(slide, 0.8, 0.18, 11.5, 1.1, data.get("title", "핵심 요약"), 34,
             bold=True, color=_rgb(theme, "bg"), align=PP_ALIGN.LEFT, theme=theme)

    bullets = data.get("bullets", [])
    if not bullets:
        return
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(11.5), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:3]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(24)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size = Pt(24)
        run.font.color.rgb = _rgb(theme, "body")


def _chart_slide(prs, data: dict, theme: dict):
    categories = data.get("categories") or []
    series = data.get("series") or []
    if not categories or not series:
        _content_slide(prs, data, theme)
        return

    slide = _blank(prs, theme)
    _textbox(slide, 0.8, 0.4, 11.5, 0.9, data.get("title", ""), 32,
             bold=True, theme=theme)
    _accent_bar(slide, 1.25, theme)

    cd = ChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(s["name"], s["values"])

    chart_type = _CHART_TYPE_MAP.get(data.get("chart_type", "bar"), XL_CHART_TYPE.BAR_CLUSTERED)
    chart_shape = slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1.6), Inches(11), Inches(5.5), cd)
    chart_shape.chart.has_legend = len(series) > 1
