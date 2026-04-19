import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def build_pptx_from_images(image_paths: list[str], title: str, output_dir: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
        )

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r"[^\w가-힣]", "_", title or "presentation")
    out_path = str(Path(output_dir) / f"{safe_name}_{ts}.pptx")
    prs.save(out_path)
    return out_path
